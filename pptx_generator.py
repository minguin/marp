from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import numpy as np
import io

class ModernPresentation:
    # モダンなカラーパレット
    PRIMARY_COLOR = RGBColor(41, 105, 176)  # 濃い青
    SECONDARY_COLOR = RGBColor(83, 129, 186)  # 薄い青
    ACCENT1_COLOR = RGBColor(244, 124, 57)  # オレンジ
    ACCENT2_COLOR = RGBColor(67, 170, 139)  # 緑
    ACCENT3_COLOR = RGBColor(144, 97, 183)  # 紫
    ACCENT4_COLOR = RGBColor(246, 196, 84)  # 黄色
    ACCENT5_COLOR = RGBColor(218, 73, 91)  # 赤
    BACKGROUND_COLOR = RGBColor(248, 249, 250)  # 背景色（薄いグレー）
    TEXT_COLOR = RGBColor(33, 37, 41)  # テキスト色（濃いグレー）
    
    # 各カテゴリの色を定義
    CATEGORY_COLORS = {
        "政治": RGBColor(41, 105, 176),
        "経済": RGBColor(67, 170, 139),
        "国際": RGBColor(144, 97, 183),
        "エンターテインメント": RGBColor(244, 124, 57),
        "テクノロジー": RGBColor(83, 129, 186),
        "スポーツ": RGBColor(218, 73, 91),
        "その他": RGBColor(150, 150, 150)
    }
    
    def __init__(self):
        self.prs = Presentation()
        self.set_slide_size(widescreen=True)
        
    def set_slide_size(self, widescreen=True):
        """スライドのサイズを設定"""
        if widescreen:
            self.prs.slide_width = Inches(13.33)
            self.prs.slide_height = Inches(7.5)
    
    def add_title_slide(self, title, subtitle):
        """タイトルスライドを追加"""
        slide_layout = self.prs.slide_layouts[0]  # タイトルスライドレイアウト
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景を設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.BACKGROUND_COLOR
        
        # タイトルとサブタイトルを設定
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle
        
        # タイトルのフォーマット設定
        title_text_frame = title_shape.text_frame
        title_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_text_frame.paragraphs[0].font.size = Pt(44)
        title_text_frame.paragraphs[0].font.bold = True
        title_text_frame.paragraphs[0].font.color.rgb = self.PRIMARY_COLOR
        
        # サブタイトルのフォーマット設定
        subtitle_text_frame = subtitle_shape.text_frame
        subtitle_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_text_frame.paragraphs[0].font.size = Pt(20)
        subtitle_text_frame.paragraphs[0].font.color.rgb = self.SECONDARY_COLOR
        
        # モダンな装飾を追加
        self._add_decorative_elements(slide)
        
        return slide
    
    def add_section_slide(self, title, section_number):
        """セクションスライドを追加"""
        slide_layout = self.prs.slide_layouts[2]  # セクションタイトルスライドレイアウト
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景を設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.PRIMARY_COLOR
        
        # タイトルを設定
        title_shape = slide.shapes.title
        title_shape.text = f"{section_number}. {title}"
        
        # タイトルのフォーマット設定
        title_text_frame = title_shape.text_frame
        title_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_text_frame.paragraphs[0].font.size = Pt(40)
        title_text_frame.paragraphs[0].font.bold = True
        title_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 白色
        
        # 装飾的な要素を追加
        left = Inches(1)
        top = Inches(3)
        width = Inches(11.33)
        height = Inches(0.1)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色
        shape.line.fill.background()
        
        return slide
    
    def add_content_slide(self, title, content_list, section_number=None):
        """コンテンツスライドを追加"""
        slide_layout = self.prs.slide_layouts[1]  # タイトルとコンテンツのレイアウト
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景を設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.BACKGROUND_COLOR
        
        # タイトルを設定
        title_shape = slide.shapes.title
        if section_number:
            title_shape.text = f"{section_number}. {title}"
        else:
            title_shape.text = title
        
        # タイトルのフォーマット設定
        title_text_frame = title_shape.text_frame
        title_text_frame.paragraphs[0].font.size = Pt(32)
        title_text_frame.paragraphs[0].font.color.rgb = self.PRIMARY_COLOR
        
        # 左側のカラーバーを追加
        left = Inches(0.2)
        top = Inches(0.2)
        width = Inches(0.1)
        height = Inches(7.1)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.PRIMARY_COLOR
        shape.line.fill.background()
        
        # コンテンツリストを追加
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()  # デフォルトのテキストをクリア
        
        for item in content_list:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(12)
        
        return slide
    
    def add_two_column_slide(self, title, left_content, right_content, section_number=None):
        """2列レイアウトのスライドを追加"""
        slide_layout = self.prs.slide_layouts[3]  # 2列のレイアウト
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景を設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.BACKGROUND_COLOR
        
        # タイトルを設定
        title_shape = slide.shapes.title
        if section_number:
            title_shape.text = f"{section_number}. {title}"
        else:
            title_shape.text = title
        
        # タイトルのフォーマット設定
        title_text_frame = title_shape.text_frame
        title_text_frame.paragraphs[0].font.size = Pt(32)
        title_text_frame.paragraphs[0].font.color.rgb = self.PRIMARY_COLOR
        
        # 左側の内容を追加
        left_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
        left_text_frame = left_shape.text_frame
        
        for item in left_content:
            p = left_text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(12)
        
        # 右側の内容を追加
        right_shape = slide.shapes.add_textbox(Inches(6.83), Inches(1.5), Inches(6), Inches(5.5))
        right_text_frame = right_shape.text_frame
        
        for item in right_content:
            p = right_text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(12)
        
        # 中央に区切り線を追加
        left = Inches(6.67)
        top = Inches(1.5)
        width = Inches(0.03)
        height = Inches(5.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.SECONDARY_COLOR
        shape.line.fill.background()
        
        return slide
    
    def add_chart_slide(self, title, chart_type, chart_data, section_number=None, subtitle=None):
        """チャートを含むスライドを追加"""
        slide_layout = self.prs.slide_layouts[6]  # 白紙のレイアウト
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景を設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.BACKGROUND_COLOR
        
        # タイトルを追加
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.33)
        height = Inches(0.8)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_text_frame = title_shape.text_frame
        
        p = title_text_frame.add_paragraph()
        if section_number:
            p.text = f"{section_number}. {title}"
        else:
            p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = self.PRIMARY_COLOR
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        # サブタイトルを追加（ある場合）
        if subtitle:
            left = Inches(0.5)
            top = Inches(1)
            width = Inches(12.33)
            height = Inches(0.5)
            subtitle_shape = slide.shapes.add_textbox(left, top, width, height)
            subtitle_text_frame = subtitle_shape.text_frame
            
            p = subtitle_text_frame.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(18)
            p.font.color.rgb = self.SECONDARY_COLOR
            p.alignment = PP_ALIGN.LEFT
        
        # チャートを作成して追加
        chart_img = self._create_chart(chart_type, chart_data)
        left = Inches(1.5)
        top = Inches(1.5)
        width = Inches(10)
        height = Inches(5.5)
        slide.shapes.add_picture(chart_img, left, top, width, height)
        
        return slide
    
    def add_process_diagram_slide(self, title, process_steps, section_number=None):
        """プロセス図を含むスライドを追加"""
        slide_layout = self.prs.slide_layouts[6]  # 白紙のレイアウト
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景を設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.BACKGROUND_COLOR
        
        # タイトルを追加
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.33)
        height = Inches(0.8)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_text_frame = title_shape.text_frame
        
        p = title_text_frame.add_paragraph()
        if section_number:
            p.text = f"{section_number}. {title}"
        else:
            p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = self.PRIMARY_COLOR
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        # プロセス図を追加
        steps = len(process_steps)
        start_x = Inches(1)
        start_y = Inches(2)
        box_width = Inches(2)
        box_height = Inches(1)
        arrow_length = Inches(1)
        
        for i, step in enumerate(process_steps):
            # ボックスを追加
            box_left = start_x + i * (box_width + arrow_length)
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_left, start_y, box_width, box_height)
            
            # 交互に色を変える
            if i % 2 == 0:
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.PRIMARY_COLOR
                text_color = RGBColor(255, 255, 255)  # 白色
            else:
                shape.fill.solid()
                shape.fill.fore_color.rgb = self.SECONDARY_COLOR
                text_color = RGBColor(255, 255, 255)  # 白色
            
            # テキストを追加
            text_frame = shape.text_frame
            text_frame.text = f"{i+1}. {step}"
            text_frame.paragraphs[0].font.color.rgb = text_color
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            
            # 矢印を追加（最後のステップ以外）
            if i < steps - 1:
                arrow_left = box_left + box_width
                arrow_top = start_y + (box_height / 2) - Inches(0.1)
                arrow_shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_left, arrow_top, arrow_length, Inches(0.2))
                arrow_shape.fill.solid()
                arrow_shape.fill.fore_color.rgb = self.ACCENT1_COLOR
                arrow_shape.line.fill.background()
        
        # 説明テキストを追加
        desc_left = Inches(1)
        desc_top = Inches(3.5)
        desc_width = Inches(11.33)
        desc_height = Inches(3.5)
        desc_shape = slide.shapes.add_textbox(desc_left, desc_top, desc_width, desc_height)
        desc_text_frame = desc_shape.text_frame
        
        for i, step in enumerate(process_steps):
            p = desc_text_frame.add_paragraph()
            p.text = f"ステップ {i+1}: {step} - " + self._get_step_description(i+1)
            p.font.size = Pt(16)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(10)
        
        return slide
    
    def add_comparison_slide(self, title, before_data, after_data, section_number=None):
        """ビフォー・アフターの比較スライドを追加"""
        slide_layout = self.prs.slide_layouts[6]  # 白紙のレイアウト
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景を設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.BACKGROUND_COLOR
        
        # タイトルを追加
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.33)
        height = Inches(0.8)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_text_frame = title_shape.text_frame
        
        p = title_text_frame.add_paragraph()
        if section_number:
            p.text = f"{section_number}. {title}"
        else:
            p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = self.PRIMARY_COLOR
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        # 「BEFORE」セクション
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(6)
        height = Inches(0.5)
        before_title = slide.shapes.add_textbox(left, top, width, height)
        before_tf = before_title.text_frame
        p = before_tf.add_paragraph()
        p.text = "処理前"
        p.font.size = Pt(24)
        p.font.color.rgb = self.ACCENT5_COLOR
        p.font.bold = True
        
        # 「AFTER」セクション
        left = Inches(6.83)
        top = Inches(1.5)
        width = Inches(6)
        height = Inches(0.5)
        after_title = slide.shapes.add_textbox(left, top, width, height)
        after_tf = after_title.text_frame
        p = after_tf.add_paragraph()
        p.text = "処理後"
        p.font.size = Pt(24)
        p.font.color.rgb = self.ACCENT2_COLOR
        p.font.bold = True
        
        # BEFOREコンテンツを追加
        left = Inches(0.5)
        top = Inches(2.2)
        width = Inches(6)
        height = Inches(4.5)
        before_content = slide.shapes.add_textbox(left, top, width, height)
        before_tf = before_content.text_frame
        
        for item in before_data:
            p = before_tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(10)
        
        # AFTERコンテンツを追加
        left = Inches(6.83)
        top = Inches(2.2)
        width = Inches(6)
        height = Inches(4.5)
        after_content = slide.shapes.add_textbox(left, top, width, height)
        after_tf = after_content.text_frame
        
        for item in after_data:
            p = after_tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(10)
        
        # 中央に区切り線と矢印を追加
        # 区切り線
        left = Inches(6.67)
        top = Inches(1.5)
        width = Inches(0.03)
        height = Inches(5.5)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        line.fill.solid()
        line.fill.fore_color.rgb = self.SECONDARY_COLOR
        line.line.fill.background()
        
        # 矢印
        left = Inches(5.67)
        top = Inches(4)
        width = Inches(2)
        height = Inches(0.5)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = self.ACCENT1_COLOR
        arrow.line.fill.background()
        
        return slide
    
    def add_system_architecture_slide(self, title, section_number=None):
        """システムアーキテクチャを表示するスライドを追加"""
        slide_layout = self.prs.slide_layouts[6]  # 白紙のレイアウト
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景を設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.BACKGROUND_COLOR
        
        # タイトルを追加
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.33)
        height = Inches(0.8)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_text_frame = title_shape.text_frame
        
        p = title_text_frame.add_paragraph()
        if section_number:
            p.text = f"{section_number}. {title}"
        else:
            p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = self.PRIMARY_COLOR
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        # システムアーキテクチャの図を作成
        # データソース
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(2)
        height = Inches(1)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.ACCENT1_COLOR
        
        text_frame = shape.text_frame
        text_frame.text = "データソース\n(GoogleNews, RSS, etc.)"
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # 重複削除処理
        left = Inches(4)
        top = Inches(1.5)
        width = Inches(2)
        height = Inches(1)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.PRIMARY_COLOR
        
        text_frame = shape.text_frame
        text_frame.text = "重複削除処理"
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # カテゴリ付与
        left = Inches(7)
        top = Inches(1.5)
        width = Inches(2)
        height = Inches(1)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.PRIMARY_COLOR
        
        text_frame = shape.text_frame
        text_frame.text = "カテゴリ付与"
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # 最終フィルタリング
        left = Inches(10)
        top = Inches(1.5)
        width = Inches(2)
        height = Inches(1)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.PRIMARY_COLOR
        
        text_frame = shape.text_frame
        text_frame.text = "ユーザーベース\nフィルタリング"
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # データベース（長方形で代用）
        left = Inches(5.5)
        top = Inches(3.5)
        width = Inches(2)
        height = Inches(1)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.ACCENT2_COLOR
        
        text_frame = shape.text_frame
        text_frame.text = "データベース"
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # ユーザープロフィール（雲の代わりに丸角四角形を使用）
        left = Inches(10)
        top = Inches(3.5)
        width = Inches(2)
        height = Inches(1)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.ACCENT3_COLOR
        
        text_frame = shape.text_frame
        text_frame.text = "ユーザープロフィール"
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # フロントエンド
        left = Inches(5.5)
        top = Inches(5.5)
        width = Inches(2)
        height = Inches(1)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.ACCENT4_COLOR
        
        text_frame = shape.text_frame
        text_frame.text = "ユーザーインターフェース"
        text_frame.paragraphs[0].font.color.rgb = RGBColor(33, 37, 41)
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        
        # 矢印を追加
        self._add_arrow(slide, Inches(3), Inches(2), Inches(1), Inches(0.2))  # データソース → 重複削除
        self._add_arrow(slide, Inches(6), Inches(2), Inches(1), Inches(0.2))  # 重複削除 → カテゴリ付与
        self._add_arrow(slide, Inches(9), Inches(2), Inches(1), Inches(0.2))  # カテゴリ付与 → フィルタリング
        self._add_arrow(slide, Inches(6.5), Inches(2.5), Inches(0.2), Inches(1))  # カテゴリ付与 → データベース
        self._add_arrow(slide, Inches(10), Inches(2.5), Inches(0.2), Inches(1), is_down=False)  # フィルタリング → ユーザープロフィール
        self._add_arrow(slide, Inches(6.5), Inches(4.5), Inches(0.2), Inches(1))  # データベース → フロントエンド
        self._add_arrow(slide, Inches(10), Inches(4.5), Inches(0.2), Inches(1), is_down=True, is_horizontal=False)  # ユーザープロフィール → フロントエンド
        
        # 説明テキストを追加
        left = Inches(0.5)
        top = Inches(6.7)
        width = Inches(12.33)
        height = Inches(0.7)
        desc_shape = slide.shapes.add_textbox(left, top, width, height)
        desc_text_frame = desc_shape.text_frame
        
        p = desc_text_frame.add_paragraph()
        p.text = "このシステムは、複数のデータソースから情報を取得し、重複を除去し、カテゴリを付与した後、ユーザーの興味に基づいてフィルタリングを行います。処理されたデータはデータベースに保存され、ユーザーインターフェースを通じて提供されます。"
        p.font.size = Pt(14)
        p.font.color.rgb = self.TEXT_COLOR
        p.alignment = PP_ALIGN.LEFT
        
        return slide
    
    def add_category_chart_slide(self, title, section_number=None):
        """カテゴリ分布を表示するチャートスライドを追加"""
        slide_layout = self.prs.slide_layouts[6]  # 白紙のレイアウト
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景を設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.BACKGROUND_COLOR
        
        # タイトルを追加
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.33)
        height = Inches(0.8)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_text_frame = title_shape.text_frame
        
        p = title_text_frame.add_paragraph()
        if section_number:
            p.text = f"{section_number}. {title}"
        else:
            p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = self.PRIMARY_COLOR
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        # カテゴリ分布の円グラフを作成
        plt.figure(figsize=(10, 6))
        categories = ["政治", "経済", "国際", "エンターテインメント", "テクノロジー", "スポーツ", "その他"]
        sizes = [25, 20, 15, 15, 15, 5, 5]
        colors = [
            (41/255, 105/255, 176/255),
            (67/255, 170/255, 139/255),
            (144/255, 97/255, 183/255),
            (244/255, 124/255, 57/255),
            (83/255, 129/255, 186/255),
            (218/255, 73/255, 91/255),
            (150/255, 150/255, 150/255)
        ]
        
        plt.pie(sizes, labels=categories, colors=colors, autopct='%1.1f%%',
                shadow=False, startangle=90, textprops={'fontsize': 12})
        plt.axis('equal')
        plt.title('カテゴリ分布', fontsize=18)
        
        # 画像をバイトストリームに保存して追加
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=100)  # bbox_inches='tight' パラメータを削除
        img_stream.seek(0)
        plt.close()
        
        # 画像を追加
        left = Inches(1.5)
        top = Inches(1.5)
        width = Inches(10)
        height = Inches(5)
        slide.shapes.add_picture(img_stream, left, top, width, height)
        
        # 凡例を追加
        left = Inches(1)
        top = Inches(6.5)
        width = Inches(11.33)
        height = Inches(0.8)
        legend_shape = slide.shapes.add_textbox(left, top, width, height)
        legend_text_frame = legend_shape.text_frame
        
        p = legend_text_frame.add_paragraph()
        p.text = "カテゴリ分類の基準: ヘッドラインと記事の内容から機械学習アルゴリズムを用いて自動分類しています。ユーザーごとの関心領域に基づいて、コンテンツの重み付けを行います。"
        p.font.size = Pt(14)
        p.font.color.rgb = self.TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_table_slide(self, title, headers, rows, section_number=None):
        """表を含むスライドを追加"""
        slide_layout = self.prs.slide_layouts[6]  # 白紙のレイアウト
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景を設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.BACKGROUND_COLOR
        
        # タイトルを追加
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.33)
        height = Inches(0.8)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_text_frame = title_shape.text_frame
        
        p = title_text_frame.add_paragraph()
        if section_number:
            p.text = f"{section_number}. {title}"
        else:
            p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = self.PRIMARY_COLOR
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        # 表を追加
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(11.33)
        height = Inches(0.5)
        rows_count = len(rows) + 1  # ヘッダー行 + データ行
        cols_count = len(headers)
        
        table = slide.shapes.add_table(rows_count, cols_count, left, top, width, height).table
        
        # ヘッダー行を設定
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.PRIMARY_COLOR
            
            # ヘッダーテキストのフォーマット
            text_frame = cell.text_frame
            paragraph = text_frame.paragraphs[0]
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
        
        # データ行を設定
        for i, row_data in enumerate(rows):
            for j, cell_data in enumerate(row_data):
                cell = table.cell(i + 1, j)
                cell.text = str(cell_data)
                
                # 偶数行と奇数行で背景色を変える
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                # テキストのフォーマット
                text_frame = cell.text_frame
                paragraph = text_frame.paragraphs[0]
                paragraph.font.color.rgb = self.TEXT_COLOR
                paragraph.font.size = Pt(14)
                
                # 列によって表示形式を変更
                if j == 0:  # 最初の列は左揃え
                    paragraph.alignment = PP_ALIGN.LEFT
                else:  # その他の列は中央揃え
                    paragraph.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_summary_slide(self, title, summary_points, future_prospects, section_number=None):
        """まとめスライドを追加"""
        slide_layout = self.prs.slide_layouts[6]  # 白紙のレイアウト
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景を設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.BACKGROUND_COLOR
        
        # タイトルを追加
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(12.33)
        height = Inches(0.8)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
        title_text_frame = title_shape.text_frame
        
        p = title_text_frame.add_paragraph()
        if section_number:
            p.text = f"{section_number}. {title}"
        else:
            p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = self.PRIMARY_COLOR
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        # 「まとめ」セクション
        left = Inches(0.5)
        top = Inches(1.3)
        width = Inches(6)
        height = Inches(0.5)
        summary_title = slide.shapes.add_textbox(left, top, width, height)
        summary_tf = summary_title.text_frame
        p = summary_tf.add_paragraph()
        p.text = "プロジェクトの成果"
        p.font.size = Pt(24)
        p.font.color.rgb = self.ACCENT2_COLOR
        p.font.bold = True
        
        # まとめ内容
        left = Inches(0.5)
        top = Inches(1.9)
        width = Inches(6)
        height = Inches(5)
        summary_content = slide.shapes.add_textbox(left, top, width, height)
        summary_tf = summary_content.text_frame
        
        for point in summary_points:
            p = summary_tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(12)
        
        # 「今後の展望」セクション
        left = Inches(6.83)
        top = Inches(1.3)
        width = Inches(6)
        height = Inches(0.5)
        future_title = slide.shapes.add_textbox(left, top, width, height)
        future_tf = future_title.text_frame
        p = future_tf.add_paragraph()
        p.text = "今後の展望"
        p.font.size = Pt(24)
        p.font.color.rgb = self.ACCENT1_COLOR
        p.font.bold = True
        
        # 今後の展望内容
        left = Inches(6.83)
        top = Inches(1.9)
        width = Inches(6)
        height = Inches(5)
        future_content = slide.shapes.add_textbox(left, top, width, height)
        future_tf = future_content.text_frame
        
        for prospect in future_prospects:
            p = future_tf.add_paragraph()
            p.text = f"• {prospect}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_after = Pt(12)
        
        # 中央に区切り線を追加
        left = Inches(6.67)
        top = Inches(1.3)
        width = Inches(0.03)
        height = Inches(5.6)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        line.fill.solid()
        line.fill.fore_color.rgb = self.SECONDARY_COLOR
        line.line.fill.background()
        
        return slide
    
    def add_end_slide(self):
        """終了スライドを追加"""
        slide_layout = self.prs.slide_layouts[6]  # 白紙のレイアウト
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 背景を設定
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.PRIMARY_COLOR
        
        # 「ありがとうございました」テキストを追加
        left = Inches(0.5)
        top = Inches(3)
        width = Inches(12.33)
        height = Inches(1.5)
        thank_you = slide.shapes.add_textbox(left, top, width, height)
        thank_you_tf = thank_you.text_frame
        p = thank_you_tf.add_paragraph()
        p.text = "ご清聴ありがとうございました"
        p.font.size = Pt(44)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # お問い合わせ情報を追加
        left = Inches(0.5)
        top = Inches(5)
        width = Inches(12.33)
        height = Inches(1)
        contact = slide.shapes.add_textbox(left, top, width, height)
        contact_tf = contact.text_frame
        p = contact_tf.add_paragraph()
        p.text = "お問い合わせ: info@example.com"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _create_chart(self, chart_type, data):
        """チャートを作成してバイトストリームを返す"""
        plt.figure(figsize=(10, 6))
        
        if chart_type == "bar":
            # 棒グラフ
            categories = data.get("categories", [])
            values = data.get("values", [])
            colors = data.get("colors", ["#2969B0"] * len(categories))
            
            plt.bar(categories, values, color=colors)
            plt.xlabel(data.get("xlabel", "カテゴリ"))
            plt.ylabel(data.get("ylabel", "値"))
            plt.title(data.get("title", ""))
            plt.grid(axis='y', linestyle='--', alpha=0.7)
            plt.tight_layout()
            
        elif chart_type == "line":
            # 折れ線グラフ
            x = data.get("x", [])
            y = data.get("y", [])
            
            plt.plot(x, y, marker='o', linestyle='-', color="#2969B0", linewidth=2)
            plt.xlabel(data.get("xlabel", "X軸"))
            plt.ylabel(data.get("ylabel", "Y軸"))
            plt.title(data.get("title", ""))
            plt.grid(True, linestyle='--', alpha=0.7)
            plt.tight_layout()
            
        elif chart_type == "pie":
            # 円グラフ
            labels = data.get("labels", [])
            sizes = data.get("sizes", [])
            colors = data.get("colors", None)
            
            plt.pie(sizes, labels=labels, autopct='%1.1f%%', shadow=False, startangle=90, colors=colors)
            plt.axis('equal')
            plt.title(data.get("title", ""))
            
        elif chart_type == "scatter":
            # 散布図
            x = data.get("x", [])
            y = data.get("y", [])
            
            plt.scatter(x, y, color="#2969B0", alpha=0.6)
            plt.xlabel(data.get("xlabel", "X軸"))
            plt.ylabel(data.get("ylabel", "Y軸"))
            plt.title(data.get("title", ""))
            plt.grid(True, linestyle='--', alpha=0.7)
            plt.tight_layout()
            
        elif chart_type == "stacked_bar":
            # 積み上げ棒グラフ
            categories = data.get("categories", [])
            series = data.get("series", {})
            
            bottom = np.zeros(len(categories))
            for label, values in series.items():
                plt.bar(categories, values, bottom=bottom, label=label)
                bottom += values
            
            plt.xlabel(data.get("xlabel", "カテゴリ"))
            plt.ylabel(data.get("ylabel", "値"))
            plt.title(data.get("title", ""))
            plt.legend()
            plt.grid(axis='y', linestyle='--', alpha=0.7)
            plt.tight_layout()
        
        # 画像をバイトストリームに保存
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', dpi=100)  # bbox_inches='tight' を削除
        img_stream.seek(0)
        plt.close()
        
        return img_stream
    
    def _add_arrow(self, slide, left, top, width, height, is_down=False, is_horizontal=True):
        """矢印を追加するヘルパーメソッド"""
        if is_horizontal:
            shape_type = MSO_SHAPE.RIGHT_ARROW
        else:
            if is_down:
                shape_type = MSO_SHAPE.DOWN_ARROW
            else:
                shape_type = MSO_SHAPE.UP_ARROW
        
        arrow = slide.shapes.add_shape(shape_type, left, top, width, height)
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = self.ACCENT1_COLOR
        arrow.line.fill.background()
        
        return arrow
    
    def _add_decorative_elements(self, slide):
        """スライドに装飾的な要素を追加するヘルパーメソッド"""
        # 左下の装飾
        left = Inches(0)
        top = Inches(6.5)
        width = Inches(2)
        height = Inches(1)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.PRIMARY_COLOR
        shape.line.fill.background()
        
        # 右上の装飾
        left = Inches(11.33)
        top = Inches(0)
        width = Inches(2)
        height = Inches(1)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.ACCENT1_COLOR
        shape.line.fill.background()
        
        return
    
    def _get_step_description(self, step_number):
        """ステップ番号に応じた説明テキストを返す"""
        descriptions = {
            1: "データソースからGoogleニュースやRSSフィードを使って情報を収集します。",
            2: "収集したヘッドラインの重複を高精度で削除し、情報の冗長性を排除します。",
            3: "ヘッドラインの内容に基づいて、自動的にカテゴリを付与します。",
            4: "ユーザーの閲覧履歴やお気に入り情報に基づいて、パーソナライズドフィルタリングを行います。",
            5: "フィルタリングされた情報をユーザーに提供します。"
        }
        
        return descriptions.get(step_number, "このステップでは処理を行います。")


def create_news_processing_presentation():
    """ニュース処理システムのプレゼンテーションを作成する"""
    prs = ModernPresentation()
    
    # 1. タイトルスライド
    prs.add_title_slide(
        "ニュース処理システム",
        "重複削除、カテゴリ付与、パーソナライズドフィルタリング"
    )
    
    # 2. はじめに
    intro_content = [
        "プロジェクトの目的: 複数のニュースソースから収集した情報を整理し、ユーザーに最適な情報を提供する",
        "現状の問題点: 情報ソースが玉石混交しており、重複や不要な情報が多い",
        "GoogleNews等、複数の情報ソースの活用: 多様な視点からニュースを取得し、より包括的な情報を提供",
        "情報過多の時代におけるキュレーションの重要性: ユーザーに必要な情報を効率的に届ける"
    ]
    prs.add_content_slide("はじめに", intro_content, "1")
    
    # 3. 全体フローの概要
    prs.add_process_diagram_slide(
        "全体フローの概要",
        ["データ収集", "重複削除", "カテゴリ付与", "ユーザーベースフィルタリング", "コンテンツ提供"],
        "2"
    )
    
    # 4. ヘッドラインの重複削除
    duplicate_content = [
        "重複判定の基準:",
        "- タイトルの類似度（テキスト一致率70%以上を重複と判定）",
        "- 発行日時（24時間以内の類似記事を検出）",
        "- ソースの違いを考慮（同一ソースからの記事は優先的に重複除外）",
        "",
        "使用アルゴリズム:",
        "- テキスト正規化（大文字小文字の統一、特殊文字の除去）",
        "- シソーラス利用（同義語を同一視）",
        "- TF-IDFベクトル化と余弦類似度計算"
    ]
    
    before_examples = [
        "東京オリンピック開催が決定、IOCが正式発表",
        "IOC、東京オリンピックの開催を正式に発表",
        "東京五輪の開催がIOCにより正式決定",
        "急速に発達する台風19号、週末に関東直撃の恐れ",
        "台風19号、週末に関東地方へ接近する見込み",
        "新型スマートフォンの販売開始、初日から行列"
    ]
    
    after_examples = [
        "東京オリンピック開催が決定、IOCが正式発表",
        "急速に発達する台風19号、週末に関東直撃の恐れ",
        "新型スマートフォンの販売開始、初日から行列"
    ]
    
    prs.add_content_slide("ヘッドラインの重複削除", duplicate_content, "3.1")
    prs.add_comparison_slide("重複削除の処理例", before_examples, after_examples, "3.1")
    
    # 5. ニュースの分類・カテゴリ付与
    category_content = [
        "ヘッドラインを基にした自動分類ルール:",
        "- キーワードベースの初期分類",
        "- 機械学習モデルによる多クラス分類（BERT事前学習モデル活用）",
        "- コンテキスト解析による精度向上",
        "",
        "カテゴリ例:",
        "- 政治: 国内政治、選挙、法案、政策など",
        "- 経済: 株式市場、企業活動、金融、雇用など",
        "- 国際: 海外情勢、国際関係、外交など",
        "- エンターテインメント: 芸能、映画、音楽など",
        "- テクノロジー: IT、科学技術、イノベーションなど",
        "- スポーツ: 各種スポーツニュース、試合結果など"
    ]
    prs.add_content_slide("ニュースの分類・カテゴリ付与", category_content, "3.2")
    prs.add_category_chart_slide("カテゴリ分布の例", "3.2")
    
    # 6. ユーザー行動に基づく最終フィルタリング
    filter_content = [
        "閲覧履歴からのパーソナライズ抽出方法:",
        "- ユーザーの過去90日間の閲覧パターン分析",
        "- クリック率、滞在時間等のエンゲージメント指標を反映",
        "- 協調フィルタリングによる類似ユーザーからの推薦",
        "",
        "フィルタリング条件の例:",
        "- 興味カテゴリへの重み付け（閲覧頻度に基づく）",
        "- 時間経過による関心度の減衰モデル適用",
        "- 明示的なユーザー設定（フォロー、ブロック機能）",
        "",
        "ユーザー体験向上のためのフィードバックループ:",
        "- 明示的フィードバック（「いいね」「興味なし」ボタン）",
        "- 暗黙的フィードバック（閲覧時間、スクロール速度）",
        "- A/Bテストによるアルゴリズム改善"
    ]
    prs.add_content_slide("ユーザー行動に基づく最終フィルタリング", filter_content, "3.3")
    
    # フィルタリング効果のグラフ
    chart_data = {
        "categories": ["政治", "経済", "国際", "エンターテインメント", "テクノロジー", "スポーツ", "その他"],
        "series": {
            "フィルタリング前": [25, 20, 15, 15, 15, 5, 5],
            "フィルタリング後": [10, 30, 5, 30, 20, 5, 0]
        },
        "xlabel": "カテゴリ",
        "ylabel": "記事の割合（%）",
        "title": "ユーザーフィルタリング前後の記事分布変化"
    }
    prs.add_chart_slide("フィルタリング効果", "stacked_bar", chart_data, "3.3")
    
    # 7. システム全体のアーキテクチャ
    prs.add_system_architecture_slide("システム全体のアーキテクチャ", "4")
    
    # 8. 実装上のポイントと留意点
    implementation_left = [
        "精度向上のための工夫:",
        "- 言語処理前のテキスト正規化（Unicode正規化、HTML除去）",
        "- 自然言語処理ライブラリの活用（MeCab, GiNZA, spaCy）",
        "- 定期的なモデル再学習と評価指標のモニタリング"
    ]
    
    implementation_right = [
        "パフォーマンスとスケーラビリティ:",
        "- 分散処理フレームワーク（Apache Spark）の活用",
        "- キャッシュ戦略によるレスポンス時間の最適化",
        "- コンテナ化（Docker）とオートスケーリング設定",
        "",
        "セキュリティとプライバシー:",
        "- ユーザーデータの匿名化と暗号化",
        "- アクセス制御とAPI認証の厳格化",
        "- GDPRなどの規制への対応"
    ]
    
    prs.add_two_column_slide("実装上のポイントと留意点", implementation_left, implementation_right, "5")
    
    # 9. デザインおよび見せ方の提案
    design_content = [
        "ビジュアルフロー図・インフォグラフィック:",
        "- プロセスの流れを視覚的に表現し、ユーザーの理解を促進",
        "- アイコンと色を一貫して使用し、情報の関連性を強調",
        "",
        "ダッシュボードの活用:",
        "- リアルタイムデータ可視化によるシステムパフォーマンス監視",
        "- カテゴリ分布、フィルタリング効果などをインタラクティブに表示",
        "",
        "カラースキームとアイコンの統一:",
        "- 各カテゴリに一貫した色を割り当て（例: 政治=青、経済=緑）",
        "- 情報の階層構造を色の濃淡で表現し、重要度を視覚化",
        "",
        "ユーザーシナリオの事例紹介:",
        "- 具体的なユースケースを提示し、システムの有用性を実証",
        "- 実際のユーザーフィードバックを含め、説得力を高める"
    ]
    prs.add_content_slide("デザインおよび見せ方の提案", design_content, "6")
    
    # 具体的なユーザーインターフェースの例を表示
    chart_data = {
        "x": ["1月", "2月", "3月", "4月", "5月", "6月"],
        "y": [85, 88, 90, 87, 91, 95],
        "xlabel": "期間",
        "ylabel": "ユーザー満足度（%）",
        "title": "システム導入後のユーザー満足度推移"
    }
    prs.add_chart_slide("ユーザー満足度の推移", "line", chart_data, "6", "システム改善による効果測定")
    
    # 10. まとめと今後の展望
    summary_points = [
        "複数ソースからの情報統合により、情報の網羅性が30%向上",
        "重複削除アルゴリズムにより、冗長なコンテンツが75%削減",
        "自動カテゴリ付与の精度が85%に到達（人手による分類との一致率）",
        "パーソナライズドフィルタリングにより、ユーザーエンゲージメントが40%増加",
        "システム全体のスケーラビリティと安定性を確保"
    ]
    
    future_prospects = [
        "マルチモーダル分析の導入（テキスト+画像+動画）",
        "リアルタイム分析エンジンの強化によるレイテンシ削減",
        "ユーザーフィードバックループの自動化と強化学習の導入",
        "言語横断的な分析機能の拡張（多言語対応）",
        "より細粒度なカテゴリ分類と個人化アルゴリズムの高度化"
    ]
    
    prs.add_summary_slide("まとめと今後の展望", summary_points, future_prospects, "7")
    
    # 11. 終了スライド
    prs.add_end_slide()
    
    # プレゼンテーションを保存
    prs.prs.save("news_processing_system.pptx")
    return prs.prs

# プレゼンテーションを作成
if __name__ == "__main__":
    create_news_processing_presentation()
    print("プレゼンテーションが正常に作成されました。")